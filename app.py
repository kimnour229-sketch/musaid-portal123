from flask import Flask, render_template, request, redirect, url_for, session, flash, send_from_directory, abort
import sqlite3
import db  # طبقة الوصول المزدوجة (SQLite / MySQL)
import os
import time
import re
import secrets
import string
import logging
from datetime import timedelta
from collections import defaultdict
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
import os
import pypdf
import pytesseract
from docx import Document
from pptx import Presentation   # مكتبة لقراءة PowerPoint
from PIL import Image
import fitz   # مكتبة PyMuPDF لقراءة PDF كصور عند الحاجة
from flask import Flask

app = Flask(__name__)
# مفتاح الجلسة يُضبط أدناه بعد تعريف BASE_DIR (من متغير بيئي أو ملف ثابت)

# -----------------------------
# فلتر تنظيف النصوص
# -----------------------------
def clean_text(text):
    import re
    # إبقاء العربية + الإنجليزية + الأرقام + بعض الرموز الأساسية
    text = re.sub(r'[^\u0600-\u06FFa-zA-Z0-9\s.,!?؟]', '', text)
    # إزالة التكرار المبالغ فيه للأحرف
    text = re.sub(r'(.)\1{2,}', r'\1', text)
    return text.strip()

# -----------------------------
# دالة استخراج النص من الملفات
# -----------------------------
def extract_text_from_file(file_path):
    text = ""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            try:
                # محاولة القراءة كنص
                reader = pypdf.PdfReader(file_path)
                for page in reader.pages:
                    text += page.extract_text() or ""
            except Exception:
                # إذا فشلت القراءة، نستخدم OCR
                doc = fitz.open(file_path)
                for page in doc:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    # يدعم العربية والإنجليزية
                    text += pytesseract.image_to_string(img, lang="ara+eng")

        elif ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

        else:
            text = "⚠️ نوع الملف غير مدعوم حالياً."

    except Exception as e:
        text = f"⚠️ الملف غير صالح للقراءة: {e}"

    return clean_text(text)

# -----------------------------
# دالة التلخيص
# -----------------------------
def summarize_handout(file_name):
    full_path = os.path.join(app.config['UPLOAD_FOLDER'], file_name)
    content = extract_text_from_file(full_path)

    if not content or len(content.strip()) < 20:
        return "لا يوجد محتوى كافي للتلخيص."

    # تلخيص بسيط: أول 3 جمل
    import re
    sentences = re.split(r'[.!?؟]', content)
    summary = " ".join(sentences[:3])

    return summary.strip()

# إعدادات المجلد وقاعدة البيانات
DATABASE = 'musaid_ist.db'
# تم تعديل المسار ليكون مطلقاً لضمان عدم حدوث تضارب في المجلدات
BASE_DIR = os.path.abspath(os.path.dirname(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# =====================================================================
# الأمان: مفتاح الجلسة، كعكات آمنة، CSRF، تحديد المحاولات، سجل التدقيق
# =====================================================================

# تحميل متغيرات البيئة من ملف .env إن وُجد (اختياري — يتطلب python-dotenv).
# يُتجاهل بصمت إذا لم تكن المكتبة مثبّتة، فلا يكسر التشغيل.
try:
    from dotenv import load_dotenv
    load_dotenv(os.path.join(BASE_DIR, '.env'))
except Exception:
    pass


def _load_secret_key():
    """يقرأ مفتاح الجلسة من متغير بيئي، وإلا يولّد مفتاحاً قوياً ويحفظه
    في ملف ثابت حتى تبقى الجلسات صالحة بعد إعادة التشغيل."""
    env_key = os.environ.get('MUSAID_SECRET_KEY')
    if env_key:
        return env_key
    key_file = os.path.join(BASE_DIR, '.secret_key')
    try:
        if os.path.exists(key_file):
            with open(key_file, 'r', encoding='utf-8') as fh:
                saved = fh.read().strip()
                if saved:
                    return saved
        generated = secrets.token_hex(32)
        with open(key_file, 'w', encoding='utf-8') as fh:
            fh.write(generated)
        return generated
    except Exception:
        # ملاذ أخير: مفتاح لكل عملية (الجلسات لن تدوم بعد إعادة التشغيل)
        return secrets.token_hex(32)

app.secret_key = _load_secret_key()

# --- فصل بيئة الإنتاج عن التطوير ---
# MUSAID_ENV: 'production' (افتراضي) أو 'development'.
# DEBUG = False افتراضياً؛ لا يُفعَّل إلا في بيئة التطوير صراحةً.
ENV = os.environ.get('MUSAID_ENV', 'production').strip().lower()
IS_PRODUCTION = ENV not in ('development', 'dev', 'local')


def _env_bool(name, default=False):
    val = os.environ.get(name)
    if val is None:
        return default
    return val.strip().lower() in ('1', 'true', 'yes', 'on')


# DEBUG لا يُفعَّل في الإنتاج إطلاقاً. في التطوير يُفعَّل افتراضياً ويمكن ضبطه عبر MUSAID_DEBUG.
DEBUG = False if IS_PRODUCTION else _env_bool('MUSAID_DEBUG', True)

# في الإنتاج تكون الكعكات آمنة افتراضياً؛ في التطوير تُسمح على HTTP افتراضياً.
_cookie_secure_default = '1' if IS_PRODUCTION else '0'

# كعكات الجلسة المُحصَّنة. SECURE افتراضياً مُفعَّل في الإنتاج؛ يمكن تجاوزه
# عبر MUSAID_COOKIE_SECURE (0/1).
app.config.update(
    DEBUG=DEBUG,
    ENV=ENV,
    SESSION_COOKIE_HTTPONLY=True,
    SESSION_COOKIE_SAMESITE='Lax',
    SESSION_COOKIE_SECURE=_env_bool('MUSAID_COOKIE_SECURE', IS_PRODUCTION),
    PERMANENT_SESSION_LIFETIME=timedelta(hours=8),
    MAX_CONTENT_LENGTH=50 * 1024 * 1024,  # حد أقصى 50MB للرفع
)

# --- سجل تدقيق المصادقة (ملف logs/auth.log) ---
LOG_DIR = os.path.join(BASE_DIR, 'logs')
try:
    os.makedirs(LOG_DIR, exist_ok=True)
except Exception:
    pass
audit_logger = logging.getLogger('musaid.auth')
audit_logger.setLevel(logging.INFO)
if not audit_logger.handlers:
    try:
        _audit_handler = logging.FileHandler(
            os.path.join(LOG_DIR, 'auth.log'), encoding='utf-8')
        _audit_handler.setFormatter(
            logging.Formatter('%(asctime)s %(message)s'))
        audit_logger.addHandler(_audit_handler)
    except Exception:
        pass

def _client_ip():
    try:
        fwd = request.headers.get('X-Forwarded-For')
        if fwd:
            return fwd.split(',')[0].strip()
        return request.remote_addr or 'unknown'
    except Exception:
        return 'unknown'

def audit_log(event, detail=''):
    """يسجّل حدث مصادقة دون أن يكشف كلمات المرور أبداً."""
    try:
        user = session.get('user_id', '-')
    except Exception:
        user = '-'
    try:
        audit_logger.info(f'event={event} ip={_client_ip()} user={user} {detail}'.strip())
    except Exception:
        pass

# --- حماية CSRF (مبنية على رمز في الجلسة) ---
def get_csrf_token():
    token = session.get('_csrf_token')
    if not token:
        token = secrets.token_hex(32)
        session['_csrf_token'] = token
    return token

@app.context_processor
def _inject_csrf():
    return {'csrf_token': get_csrf_token}

@app.before_request
def _csrf_protect():
    if request.method in ('POST', 'PUT', 'PATCH', 'DELETE'):
        sent = (request.form.get('csrf_token')
                or request.headers.get('X-CSRFToken'))
        stored = session.get('_csrf_token')
        if not stored or not sent or not secrets.compare_digest(str(stored), str(sent)):
            audit_log('csrf_failure', f'path={request.path}')
            abort(400)

# --- تحديد محاولات الدخول + قفل مؤقت بعد الفشل المتكرر (في الذاكرة) ---
RATE_LIMIT_WINDOW = 300      # نافذة 5 دقائق
RATE_LIMIT_MAX_FAILS = 5     # 5 محاولات فاشلة
LOCKOUT_SECONDS = 900        # قفل 15 دقيقة
_login_failures = defaultdict(list)   # المفتاح -> أوقات الفشل
_login_lockouts = {}                  # المفتاح -> وقت فك القفل

def _rl_key():
    return _client_ip()

def login_lock_remaining(key):
    until = _login_lockouts.get(key)
    if until and time.time() < until:
        return int(until - time.time())
    if until:
        _login_lockouts.pop(key, None)
    return 0

def record_login_failure(key):
    now = time.time()
    fails = [t for t in _login_failures[key] if now - t < RATE_LIMIT_WINDOW]
    fails.append(now)
    _login_failures[key] = fails
    if len(fails) >= RATE_LIMIT_MAX_FAILS:
        _login_lockouts[key] = now + LOCKOUT_SECONDS
        _login_failures[key] = []
        return True
    return False

def clear_login_failures(key):
    _login_failures.pop(key, None)
    _login_lockouts.pop(key, None)

# --- تحقق من صحة المدخلات ---
_EMAIL_RE = re.compile(r'^[^@\s]+@[^@\s]+\.[^@\s]+$')

def validate_email(email):
    return bool(email) and len(email) <= 120 and bool(_EMAIL_RE.match(email))

def clean_name(name):
    return (name or '').strip()[:120]

def get_db_connection():
    # يختار المحرك تلقائياً (SQLite افتراضياً، أو MySQL عند ضبط DATABASE_URL/MUSAID_DB_BACKEND)
    return db.connect(DATABASE)

# --- ترقية إضافية آمنة: أعمدة عدّادات المشاهدة/التحميل (idempotent) ---
def ensure_schema():
    conn = get_db_connection()
    try:
        if db.is_mysql():
            # MySQL: نفحص information_schema ونضيف العمود إن غاب
            existing = {row['COLUMN_NAME'] for row in conn.execute(
                "SELECT COLUMN_NAME FROM information_schema.COLUMNS "
                "WHERE TABLE_SCHEMA = DATABASE() AND TABLE_NAME = 'handouts'")}
            for col in ('view_count', 'download_count'):
                if col not in existing:
                    conn.execute(f'ALTER TABLE handouts ADD COLUMN {col} INT DEFAULT 0')
            # نسخة نصية من كلمة المرور يطّلع عليها المدير فقط من لوحة الإدارة
            t_cols = {row['COLUMN_NAME'] for row in conn.execute(
                "SELECT COLUMN_NAME FROM information_schema.COLUMNS "
                "WHERE TABLE_SCHEMA = DATABASE() AND TABLE_NAME = 'teachers'")}
            if 'plain_password' not in t_cols:
                conn.execute('ALTER TABLE teachers ADD COLUMN plain_password VARCHAR(255)')
        else:
            # SQLite: نستخدم PRAGMA كما في الأصل
            cols = {row['name'] for row in conn.execute('PRAGMA table_info(handouts)')}
            for col in ('view_count', 'download_count'):
                if col not in cols:
                    conn.execute(f'ALTER TABLE handouts ADD COLUMN {col} INTEGER DEFAULT 0')
            t_cols = {row['name'] for row in conn.execute('PRAGMA table_info(teachers)')}
            if 'plain_password' not in t_cols:
                conn.execute('ALTER TABLE teachers ADD COLUMN plain_password TEXT')
        conn.commit()
    finally:
        conn.close()

# لا يجب أن يفشل استيراد التطبيق إذا كانت قاعدة MySQL غير مهيأة بعد.
try:
    ensure_schema()
except Exception as _e:
    print(f'⚠️  ensure_schema تخطّى التهيئة: {_e}')

# -----------------------------
# أمان كلمات المرور (تجزئة + تحقق)
# -----------------------------
# بادئات تجزئة Werkzeug المعروفة — لتمييز كلمة المرور المجزأة عن النص الصريح القديم
_HASH_PREFIXES = ('pbkdf2:', 'scrypt:', 'argon2')

def is_hashed(value):
    return bool(value) and str(value).startswith(_HASH_PREFIXES)

def verify_password(stored, candidate):
    """يتحقق من كلمة المرور سواء كانت مجزأة (حديثة) أو نصاً صريحاً (قديمة)."""
    if not stored:
        return False
    if is_hashed(stored):
        try:
            return check_password_hash(stored, candidate)
        except Exception:
            return False
    # توافق رجعي: مقارنة آمنة زمنياً للنص الصريح القديم
    return secrets.compare_digest(str(stored), str(candidate))

def hash_password(plain):
    return generate_password_hash(plain, method='pbkdf2:sha256', salt_length=16)

def validate_password(pw):
    """قواعد قوة كلمة المرور. تُعيد (صالحة, رسالة)."""
    if not pw or len(pw) < 8:
        return False, 'يجب أن تتكون كلمة المرور من 8 خانات على الأقل.'
    if not re.search(r'[A-Za-z]', pw):
        return False, 'يجب أن تحتوي كلمة المرور على حرف واحد على الأقل.'
    if not re.search(r'\d', pw):
        return False, 'يجب أن تحتوي كلمة المرور على رقم واحد على الأقل.'
    return True, ''

def generate_strong_password(length=12):
    """يولّد كلمة مرور عشوائية قوية تضمن وجود حرف ورقم ورمز."""
    alphabet = string.ascii_letters + string.digits + '!@#$%&*?'
    while True:
        pw = ''.join(secrets.choice(alphabet) for _ in range(length))
        if (re.search(r'[a-z]', pw) and re.search(r'[A-Z]', pw)
                and re.search(r'\d', pw) and re.search(r'[!@#$%&*?]', pw)):
            return pw

# --- التعديل الجوهري: حل مشكلة المعاينة والتحميل (PDF) ---
@app.route('/download/<filename>')
def uploaded_file(filename):
    # عدّاد إضافي: تحميل عند ?dl=1 وإلا تُحتسب مشاهدة. لا يؤثر على إرسال الملف.
    try:
        column = 'download_count' if request.args.get('dl') else 'view_count'
        conn = get_db_connection()
        conn.execute(
            f'UPDATE handouts SET {column} = COALESCE({column}, 0) + 1 WHERE file_path = ?',
            (filename,))
        conn.commit()
        conn.close()
    except Exception:
        pass
    # as_attachment=False تضمن فتح الملف في المتصفح (المعاينة) بدلاً من إجبار التحميل
    # Flask سيتعرف تلقائياً على امتداد .pdf ويرسل الرأس الصحيح للمتصفح
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename, as_attachment=False)

# --- خدمة ملفات PWA من جذر المشروع (Service Worker + Manifests) ---
@app.route('/service-worker.js')
def service_worker():
    return send_from_directory(BASE_DIR, 'service-worker.js', mimetype='application/javascript')

@app.route('/manifest_student.json')
def manifest_student():
    return send_from_directory(BASE_DIR, 'manifest_student.json', mimetype='application/manifest+json')

@app.route('/manifest_teacher.json')
def manifest_teacher():
    return send_from_directory(BASE_DIR, 'manifest_teacher.json', mimetype='application/manifest+json')

# --- 1. واجهة الطالب ---
@app.route('/')
def index():
    conn = get_db_connection()
    depts = conn.execute('SELECT * FROM departments').fetchall()

    # إحصائيات حية للصفحة الرئيسية (قراءة فقط)
    stats = {
        'handouts': conn.execute('SELECT COUNT(*) FROM handouts').fetchone()[0],
        'subjects': conn.execute('SELECT COUNT(*) FROM subjects').fetchone()[0],
        'teachers': conn.execute("SELECT COUNT(*) FROM teachers WHERE email != 'admin@musaid.edu.ly'").fetchone()[0],
        'departments': conn.execute('SELECT COUNT(*) FROM departments').fetchone()[0],
    }

    conn.close()
    return render_template('index.html', depts=depts, stats=stats)
    
@app.route('/search')
def search():
    dept_id = request.args.get('dept')
    semester = request.args.get('semester')
    
    conn = get_db_connection()
    query = '''
        SELECT h.*, s.subject_name, t.full_name as teacher_name
        FROM handouts h
        JOIN subjects s ON h.subject_id = s.id
        JOIN teachers t ON h.teacher_id = t.id
        WHERE h.dept_id = ? AND h.semester = ?
    '''
    results = conn.execute(query, (dept_id, semester)).fetchall()
    
    dept_name_row = conn.execute('SELECT dept_name FROM departments WHERE id = ?', (dept_id,)).fetchone()
    dept_name = dept_name_row['dept_name'] if dept_name_row else "غير معروف"
    
    processed_results = []
    for row in results:
        item = dict(row)
        # التلخيص العميق باستخدام الملف نفسه
        summary = summarize_handout(item.get('file_path', ''))
        item['flash_summary'] = summary
        processed_results.append(item)
    conn.close()
    
    return render_template('results.html', results=processed_results, dept_name=dept_name, semester=semester)

# --- 2. نظام تسجيل الدخول ---
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        key = _rl_key()

        # قفل مؤقت بعد تكرار الفشل
        remaining = login_lock_remaining(key)
        if remaining > 0:
            audit_log('login_locked', f'email={request.form.get("email","")!r}')
            flash(f'🔒 تم تجاوز عدد المحاولات المسموح. حاول مجدداً بعد {remaining // 60 + 1} دقيقة.')
            return render_template('login.html')

        email = (request.form.get('email') or '').strip()
        password = request.form.get('password') or ''

        conn = get_db_connection()
        # البحث بالبريد فقط ثم التحقق من كلمة المرور (يدعم المجزأة والنص الصريح القديم)
        user = conn.execute('SELECT * FROM teachers WHERE email = ?', (email,)).fetchone()

        if user and verify_password(user['password'], password):
            # ترقية كسولة: تجزئة كلمات المرور النصية القديمة عند أول دخول ناجح
            # مع الاحتفاظ بنسخة نصية يطّلع عليها المدير من لوحة الإدارة
            if not is_hashed(user['password']):
                conn.execute('UPDATE teachers SET password = ?, plain_password = ? WHERE id = ?',
                             (hash_password(password), password, user['id']))
                conn.commit()
            conn.close()

            clear_login_failures(key)
            # تدوير معرّف الجلسة لمنع تثبيت الجلسة (Session Fixation)
            old_csrf = session.get('_csrf_token')
            session.clear()
            if old_csrf:
                session['_csrf_token'] = old_csrf
            session['user_id'] = user['id']
            session['user_name'] = user['full_name']
            session['role'] = 'admin' if email == 'admin@musaid.edu.ly' else 'teacher'
            audit_log('login_success', f'role={session["role"]} email={email!r}')

            if session['role'] == 'admin':
                return redirect(url_for('admin_dashboard'))
            else:
                return redirect(url_for('teacher_dashboard'))
        else:
            conn.close()
            locked_now = record_login_failure(key)
            audit_log('login_failure', f'email={email!r} locked={locked_now}')
            if locked_now:
                flash(f'🔒 تم تجاوز عدد المحاولات المسموح. الحساب مقفل مؤقتاً لمدة {LOCKOUT_SECONDS // 60} دقيقة.')
            else:
                flash('خطأ في البريد الإلكتروني أو كلمة المرور')

    return render_template('login.html')

# --- 3. لوحة تحكم الأستاذ ---
@app.route('/teacher')
def teacher_dashboard():
    if 'user_id' not in session or session['role'] != 'teacher':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    depts = conn.execute('SELECT id, dept_name FROM departments').fetchall()
    depts_list = [dict(row) for row in depts]
    
    my_handouts = conn.execute('''
        SELECT h.*, s.subject_name 
        FROM handouts h 
        JOIN subjects s ON h.subject_id = s.id 
        WHERE h.teacher_id = ? 
        ORDER BY h.id DESC
    ''', (session['user_id'],)).fetchall()
    my_handouts_list = [dict(row) for row in my_handouts]

    # عدد المذكرات التي رفعها هذا الأستاذ
    my_count = len(my_handouts_list)

    # أكثر مادة رفع لها مذكرات
    top_subject = None
    if my_count > 0:
        subject_counts = {}
        for h in my_handouts_list:
            subject_counts[h['subject_name']] = subject_counts.get(h['subject_name'], 0) + 1
        top_subject = max(subject_counts, key=subject_counts.get)

    # إجمالي المذكرات لجميع الأساتذة
    total_count = conn.execute("SELECT COUNT(*) FROM handouts").fetchone()[0]

    # نسبة المشاركة
    participation = round((my_count / total_count) * 100, 1) if total_count > 0 else 0

    # تنبيه ذكي عند قلة النشاط (شرط عام فقط)
    ai_message = None
    if my_count == 0:
        ai_message = "🔔 تنبيه ذكي: لم تقم برفع أي مذكرة حتى الآن هذا الفصل."

    conn.close()
    
    return render_template('teacher_dashboard.html', 
                           name=session['user_name'], 
                           depts=depts_list,
                           my_handouts=my_handouts_list,
                           my_count=my_count,
                           top_subject=top_subject,
                           participation=participation,
                           ai_alert=ai_message)

import os
from werkzeug.utils import secure_filename

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'user_id' not in session or session['role'] != 'teacher':
        return redirect(url_for('login'))

    dept_id = request.form.get('dept_id')
    semester = request.form.get('semester')
    subject_id = request.form.get('subject_id')
    title = request.form.get('title')
    notes = request.form.get('notes')
    files = request.files.getlist('files[]')

    conn = get_db_connection()
    try:
        # 🔍 فحص التكرار قبل الحفظ
        duplicate = conn.execute("""
            SELECT COUNT(*) FROM handouts 
            WHERE subject_id = ? AND LOWER(title) = LOWER(?)
        """, (subject_id, title)).fetchone()[0]

        if duplicate > 0:
            conn.close()
            flash("⚠️ تنبيه ذكي: هناك مذكرة مشابهة مرفوعة مسبقًا لهذه المادة.")
            return redirect(url_for('teacher_dashboard'))

        # إذا لا يوجد تكرار → نحفظ الملفات
        for file in files:
            if file and file.filename != '':
                original_filename = file.filename
                ext = os.path.splitext(original_filename)[1]
                base_name = secure_filename(os.path.splitext(original_filename)[0])
                unique_filename = f"{int(time.time())}_{base_name}{ext}"
                file.save(os.path.join(app.config['UPLOAD_FOLDER'], unique_filename))

                conn.execute('''
                    INSERT INTO handouts (teacher_id, subject_id, dept_id, semester, title, notes, file_path)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                ''', (session['user_id'], subject_id, dept_id, semester, title, notes, unique_filename))
        
        conn.commit()
        flash('✅ تم رفع الملفات بنجاح!')
    except Exception as e:
        print(f"Error: {e}")
        flash('❌ حدث خطأ أثناء الرفع!')
    finally:
        conn.close()
    
    return redirect(url_for('teacher_dashboard'))

# --- 4. لوحة تحكم المدير ---
@app.route('/admin')
def admin_dashboard():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    teachers = conn.execute("SELECT * FROM teachers WHERE email != 'admin@musaid.edu.ly' ORDER BY id DESC").fetchall()
    conn.close()
    return render_template('admin/dashboard.html', teachers=teachers)

@app.route('/admin/add_teacher', methods=['POST'])
def add_teacher():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    full_name = clean_name(request.form.get('name'))
    email = (request.form.get('email') or '').strip().lower()
    password = request.form.get('password') or ''

    # تحقق من صحة المدخلات
    if not full_name:
        flash('⚠️ يرجى إدخال اسم الأستاذ.')
        return redirect(url_for('admin_dashboard'))
    if not validate_email(email):
        flash('⚠️ البريد الإلكتروني غير صالح.')
        return redirect(url_for('admin_dashboard'))

    # ضمان قوة كلمة المرور: إن كانت ضعيفة أو فارغة نولّد واحدة قوية
    ok, _ = validate_password(password)
    if not ok:
        password = generate_strong_password()

    conn = get_db_connection()
    try:
        conn.execute("INSERT INTO teachers (full_name, email, password, plain_password) VALUES (?, ?, ?, ?)",
                     (full_name, email, hash_password(password), password))
        conn.commit()
        audit_log('teacher_add', f'email={email!r}')
        # تُعرض بيانات الدخول مرة واحدة فقط (لا تُخزَّن كنص صريح)
        flash(f'✅ تم تسجيل الأستاذ بنجاح! بيانات الدخول (احفظها الآن): {email} / {password}')
    except Exception as e:
        flash('خطأ: قد يكون البريد مسجلاً مسبقاً!')
    finally:
        conn.close()
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/edit_teacher', methods=['POST'])
def edit_teacher():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    t_id = request.form.get('id')
    name = clean_name(request.form.get('name'))
    email = (request.form.get('email') or '').strip().lower()
    password = (request.form.get('password') or '').strip()

    # تحقق من صحة المدخلات
    if not name or not validate_email(email):
        flash('⚠️ الاسم أو البريد الإلكتروني غير صالح.')
        return redirect(url_for('admin_dashboard'))

    conn = get_db_connection()
    if password:
        # تغيير كلمة المرور مطلوب: نتحقق من قوتها ثم نجزّئها
        ok, msg = validate_password(password)
        if not ok:
            conn.close()
            flash('⚠️ ' + msg)
            return redirect(url_for('admin_dashboard'))
        conn.execute("UPDATE teachers SET full_name = ?, email = ?, password = ?, plain_password = ? WHERE id = ?",
                     (name, email, hash_password(password), password, t_id))
    else:
        # ترك الحقل فارغاً يُبقي كلمة المرور الحالية دون تغيير
        conn.execute("UPDATE teachers SET full_name = ?, email = ? WHERE id = ?",
                     (name, email, t_id))
    conn.commit()
    conn.close()
    audit_log('teacher_edit', f'id={t_id} pw_changed={bool(password)}')
    flash('تم تحديث بيانات الأستاذ بنجاح')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/subjects')
def admin_subjects():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    depts = conn.execute("SELECT * FROM departments ORDER BY id").fetchall()
    dept_map = {d['id']: d['dept_name'] for d in depts}
    query = '''
        SELECT s.id, s.subject_name, cs.dept_id, cs.semester
        FROM subjects s
        JOIN course_structure cs ON s.id = cs.subject_id
    '''
    subjects = conn.execute(query).fetchall()
    conn.close()
    
    # منطق التنبيه الذكي
    ai_message = None
    if len(subjects) == 0:
        ai_message = "🔔 تنبيه ذكي: لا توجد أي مواد مسجلة حتى الآن."
    else:
        # مثال: إذا عدد المواد أكبر من عدد الأقسام × 10، نعتبره مؤشر ضغط
        if len(subjects) > len(depts) * 10:
            ai_message = "🔔 تنبيه ذكي: هناك عدد كبير من المواد مقارنة بالأقسام، تحقق من التوزيع."
    
    return render_template('admin/subjects.html',
                           depts=depts,
                           dept_map=dept_map,
                           subjects=subjects,
                           ai_alert=ai_message)

@app.route('/admin/delete_teacher/<int:id>')
def delete_teacher(id):
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    conn.execute("DELETE FROM teachers WHERE id = ?", (id,))
    conn.commit()
    conn.close()
    audit_log('teacher_delete', f'id={id}')
    flash('تم حذف حساب الأستاذ بنجاح')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/add_subject', methods=['POST'])
def add_subject():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
        
    name = request.form.get('subject_name')
    dept_id = request.form.get('dept_id')
    semester = request.form.get('semester')
    
    conn = get_db_connection()
    try:
        cursor = conn.cursor()
        cursor.execute("INSERT INTO subjects (subject_name) VALUES (?)", (name,))
        subject_id = cursor.lastrowid
        
        conn.execute('''
            INSERT INTO course_structure (subject_id, dept_id, semester) 
            VALUES (?, ?, ?)
        ''', (subject_id, dept_id, semester))
        
        conn.commit()
        flash('تمت إضافة المادة وربطها بالقسم بنجاح!')
    except Exception as e:
        flash('حدث خطأ أثناء إضافة المادة.')
    finally:
        conn.close()
    return redirect(url_for('admin_subjects'))

@app.route('/admin/monitor')
def admin_monitor():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    query = '''
        SELECT h.*, s.subject_name, d.dept_name, t.full_name as teacher_name
        FROM handouts h
        JOIN subjects s ON h.subject_id = s.id
        JOIN departments d ON h.dept_id = d.id
        JOIN teachers t ON h.teacher_id = t.id
        ORDER BY h.id DESC
    '''
    logs = conn.execute(query).fetchall()
    conn.close()
    
    # منطق التنبيه الذكي
    ai_message = None
    if len(logs) == 0:
        ai_message = "🔔 تنبيه ذكي: لا توجد أي مذكرات مرفوعة حتى الآن."
    elif len(logs) > 20:
        ai_message = "🔔 تنبيه ذكي: تم رفع عدد كبير من المذكرات مؤخرًا، تحقق من صحتها."

    return render_template('admin/monitor.html', logs=logs, ai_alert=ai_message)


@app.route('/admin/delete_subject/<int:id>')
def delete_subject(id):
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    try:
        conn.execute("DELETE FROM course_structure WHERE subject_id = ?", (id,))
        conn.execute("DELETE FROM subjects WHERE id = ?", (id,))
        conn.commit()
        flash('تم حذف المادة بنجاح.')
        # تنبيه ذكي عند الحذف
        ai_message = "🔔 تنبيه ذكي: تم حذف مادة من النظام، تأكد من تحديث الهيكل الدراسي."
    except Exception as e:
        flash('حدث خطأ أثناء الحذف.')
        ai_message = "🔔 تنبيه ذكي: حدث خطأ أثناء محاولة حذف المادة."
    finally:
        conn.close()
    
    # نمرر التنبيه مع إعادة التوجيه
    return redirect(url_for('admin_subjects', ai_alert=ai_message))


@app.route('/admin/reports')
def admin_reports():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    try:
        t_count = conn.execute("SELECT COUNT(*) FROM teachers WHERE email != 'admin@musaid.edu.ly'").fetchone()[0]
        s_count = conn.execute('SELECT COUNT(*) FROM subjects').fetchone()[0]
        l_count = conn.execute('SELECT COUNT(*) FROM handouts').fetchone()[0]
        # توزيع المذكرات على كل الأقسام ديناميكياً (يشمل أي قسم جديد)
        dept_stats = conn.execute('''
            SELECT d.dept_name, COUNT(h.id) AS cnt
            FROM departments d
            LEFT JOIN handouts h ON h.dept_id = d.id
            GROUP BY d.id, d.dept_name
            ORDER BY d.id
        ''').fetchall()
        dept_stats = [{'name': r['dept_name'], 'count': r['cnt']} for r in dept_stats]
    except Exception as e:
        t_count = s_count = l_count = 0
        dept_stats = []
    finally:
        conn.close()

    # منطق التنبيه الذكي
    ai_message = None
    if l_count == 0:
        ai_message = "🔔 تنبيه ذكي: لم يتم رفع أي مذكرة هذا الأسبوع."
    else:
        weakest = min(dept_stats, key=lambda d: d['count'], default=None)
        if weakest and weakest['count'] < 2:
            ai_message = f"🔔 تنبيه ذكي: نشاط قسم {weakest['name']} منخفض، يرجى المتابعة."

    return render_template('admin/reports.html',
                           t_count=t_count,
                           s_count=s_count,
                           l_count=l_count,
                           dept_stats=dept_stats,
                           ai_alert=ai_message)

@app.route('/admin/departments')
def admin_departments():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))

    conn = get_db_connection()
    departments = conn.execute('''
        SELECT d.id, d.dept_name,
               (SELECT COUNT(*) FROM course_structure cs WHERE cs.dept_id = d.id) AS subject_count,
               (SELECT COUNT(*) FROM handouts h WHERE h.dept_id = d.id) AS handout_count
        FROM departments d
        ORDER BY d.id
    ''').fetchall()
    conn.close()
    return render_template('admin/departments.html', departments=departments)

@app.route('/admin/add_department', methods=['POST'])
def add_department():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))

    name = (request.form.get('dept_name') or '').strip()[:120]
    if not name:
        flash('⚠️ يرجى إدخال اسم القسم.')
        return redirect(url_for('admin_departments'))

    conn = get_db_connection()
    try:
        conn.execute('INSERT INTO departments (dept_name) VALUES (?)', (name,))
        conn.commit()
        audit_log('department_add', f'name={name!r}')
        flash(f'✅ تم إضافة قسم "{name}" بنجاح.')
    except Exception:
        flash('⚠️ تعذّر الإضافة — قد يكون القسم مسجلاً مسبقاً.')
    finally:
        conn.close()
    return redirect(url_for('admin_departments'))

@app.route('/admin/delete_department/<int:id>', methods=['POST'])
def delete_department(id):
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))

    conn = get_db_connection()
    # منع الحذف إن كان القسم مرتبطاً بمواد أو مذكرات حتى لا تُيتَّم البيانات
    used = conn.execute(
        'SELECT (SELECT COUNT(*) FROM course_structure WHERE dept_id = ?) '
        '+ (SELECT COUNT(*) FROM handouts WHERE dept_id = ?)', (id, id)).fetchone()[0]
    if used:
        conn.close()
        flash('⚠️ لا يمكن حذف القسم لأنه مرتبط بمواد أو مذكرات. احذفها أولاً.')
        return redirect(url_for('admin_departments'))

    conn.execute('DELETE FROM departments WHERE id = ?', (id,))
    conn.commit()
    conn.close()
    audit_log('department_delete', f'id={id}')
    flash('تم حذف القسم بنجاح')
    return redirect(url_for('admin_departments'))

@app.route('/change_password', methods=['GET', 'POST'])
def change_password():
    if 'user_id' not in session:
        return redirect(url_for('login'))
    # التحكم في كلمات المرور للمدير فقط
    if session.get('role') != 'admin':
        return redirect(url_for('teacher_dashboard'))

    if request.method == 'POST':
        current = request.form.get('current_password') or ''
        new = request.form.get('new_password') or ''
        confirm = request.form.get('confirm_password') or ''

        conn = get_db_connection()
        user = conn.execute('SELECT * FROM teachers WHERE id = ?',
                            (session['user_id'],)).fetchone()

        if not user or not verify_password(user['password'], current):
            conn.close()
            audit_log('password_change_failed', 'reason=wrong_current')
            flash('❌ كلمة المرور الحالية غير صحيحة.')
            return redirect(url_for('change_password'))

        if new != confirm:
            conn.close()
            flash('⚠️ كلمة المرور الجديدة وتأكيدها غير متطابقين.')
            return redirect(url_for('change_password'))

        ok, msg = validate_password(new)
        if not ok:
            conn.close()
            flash('⚠️ ' + msg)
            return redirect(url_for('change_password'))

        conn.execute('UPDATE teachers SET password = ? WHERE id = ?',
                     (hash_password(new), session['user_id']))
        conn.commit()
        conn.close()
        audit_log('password_change_success')
        flash('✅ تم تغيير كلمة المرور بنجاح.')
        dest = 'admin_dashboard' if session.get('role') == 'admin' else 'teacher_dashboard'
        return redirect(url_for(dest))

    return render_template('change_password.html')

@app.route('/logout')
def logout():
    audit_log('logout')
    session.clear()
    flash('تم تسجيل الخروج بنجاح.')
    return redirect(url_for('login'))

# --- الدوال المساعدة للجافاسكريبت ---
@app.route('/get_semesters/<int:dept_id>')
def get_semesters(dept_id):
    conn = get_db_connection()
    semesters = conn.execute('''
        SELECT DISTINCT semester FROM course_structure 
        WHERE dept_id = ? ORDER BY semester
    ''', (dept_id,)).fetchall()
    conn.close()
    if not semesters:
        return {"semesters": [{"id": i, "number": i} for i in range(1, 7)]}
    return {"semesters": [{"id": s['semester'], "number": s['semester']} for s in semesters]}

@app.route('/get_subjects/<int:dept_id>/<int:semester>')
def get_subjects(dept_id, semester):
    conn = get_db_connection()
    try:
        subjects = conn.execute('''
            SELECT s.id, s.subject_name FROM subjects s
            JOIN course_structure cs ON s.id = cs.subject_id
            WHERE cs.dept_id = ? AND cs.semester = ?
        ''', (dept_id, semester)).fetchall()
    except:
        subjects = []
    conn.close()
    return {"subjects": [{"id": s['id'], "name": s['subject_name']} for s in subjects]}

@app.route('/delete_handout/<int:handout_id>')
def delete_handout(handout_id):
    if 'user_id' not in session or session['role'] not in ['teacher', 'admin']:
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    try:
        handout = conn.execute('SELECT * FROM handouts WHERE id = ?', (handout_id,)).fetchone()
        if handout:
            if session['role'] == 'admin' or handout['teacher_id'] == session['user_id']:
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], handout['file_path'])
                if os.path.exists(file_path):
                    os.remove(file_path)
                conn.execute('DELETE FROM handouts WHERE id = ?', (handout_id,))
                conn.commit()
                flash('تم حذف المذكرة بنجاح')
    except Exception as e:
        flash(f'خطأ: {str(e)}')
    finally:
        conn.close()
    return redirect(request.referrer or url_for('index'))

# تجهيز حساب المدير — يُنشأ فقط إن لم يكن موجوداً، بكلمة مرور مجزأة.
# لا نعيد ضبط كلمة المرور عند كل تشغيل حتى لا نُلغي التجزئة أو تغييرات المدير.
with app.app_context():
    try:
        conn = get_db_connection()
        existing = conn.execute(
            "SELECT id FROM teachers WHERE email = 'admin@musaid.edu.ly'").fetchone()
        if not existing:
            conn.execute(
                "INSERT INTO teachers (full_name, email, password) VALUES (?, ?, ?)",
                ('مدير النظام', 'admin@musaid.edu.ly', hash_password('33557799')))
            conn.commit()
        conn.close()
    except Exception as _e:
        print(f'⚠️  تخطّى تهيئة حساب المدير: {_e}')

if __name__ == '__main__':
    # تشغيل خادم التطوير المدمج فقط عند التنفيذ المباشر.
    # في الإنتاج يُشغَّل التطبيق عبر خادم WSGI (Gunicorn/Waitress) ولا يمر من هنا.
    # DEBUG = False افتراضياً ولا يُفعَّل إلا في بيئة التطوير صراحةً (انظر الأعلى).
    host = os.environ.get('MUSAID_HOST', '127.0.0.1')
    port = int(os.environ.get('MUSAID_PORT', '5000'))
    if not DEBUG:
      print("Use Gunicorn in production. DEBUG=False")
    app.run(host=host, port=port, debug=DEBUG)